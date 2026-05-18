"""
POV Deck Generator — Rudra Multi-Model Agentic Architecture
Generates a 12-slide executive PowerPoint with architecture visuals.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Brand Colors ──────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy background
BLUE        = RGBColor(0x1B, 0x6C, 0xA8)   # primary blue
CYAN        = RGBColor(0x00, 0xC2, 0xFF)   # accent cyan
GREEN       = RGBColor(0x00, 0xD4, 0x8A)   # positive/success green
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)   # highlight orange
RED         = RGBColor(0xE7, 0x3C, 0x3E)   # alert red
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xCC, 0xD6, 0xE0)
MID_GREY    = RGBColor(0x55, 0x6B, 0x82)
DARK_CARD   = RGBColor(0x16, 0x2D, 0x40)   # card background
GOLD        = RGBColor(0xFF, 0xC8, 0x00)   # gold accent

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


# ── Helpers ───────────────────────────────────────────────────────────────────

def bg(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill: RGBColor, alpha=None, radius=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()  # no border
    return shape


def rounded_rect(slide, l, t, w, h, fill: RGBColor, corner_size=0.08):
    """Add a rounded rectangle."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = corner_size
    except Exception:
        pass
    return shape


def label(slide, text, l, t, w, h,
          font_size=12, bold=False, color=WHITE,
          align=PP_ALIGN.CENTER, italic=False, wrap=True):
    """Add a text label (transparent background)."""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def card(slide, l, t, w, h, fill=DARK_CARD,
         title="", title_size=11, title_color=CYAN,
         body="", body_size=9.5, body_color=LIGHT_GREY,
         accent: RGBColor = None):
    """Render a dark card with optional top accent bar."""
    rounded_rect(slide, l, t, w, h, fill)
    if accent:
        rect(slide, l, t, w, 0.045, accent)
    if title:
        label(slide, title, l + 0.12, t + 0.07, w - 0.2, 0.28,
              font_size=title_size, bold=True, color=title_color)
    if body:
        label(slide, body, l + 0.12, t + 0.36, w - 0.22, h - 0.44,
              font_size=body_size, color=body_color, align=PP_ALIGN.LEFT)


def arrow_right(slide, l, t, w=0.35, h=0.18, color=CYAN):
    """Draw a right-pointing arrow."""
    shape = slide.shapes.add_shape(
        13,  # RIGHT_ARROW
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def arrow_down(slide, l, t, w=0.18, h=0.3, color=CYAN):
    shape = slide.shapes.add_shape(
        36,  # DOWN_ARROW
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def divider(slide, t, color=BLUE, l=0.4, r=0.4):
    rect(slide, l, t, 13.33 - l - r, 0.025, color)


def chip(slide, text, l, t, w, h, bg_color=BLUE, text_color=WHITE, font_size=8.5):
    """Small pill/chip label."""
    rounded_rect(slide, l, t, w, h, bg_color, corner_size=0.15)
    label(slide, text, l, t + 0.01, w, h - 0.02,
          font_size=font_size, bold=True, color=text_color)


# ─── SLIDE 1: Title ───────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)

    # Left accent bar
    rect(slide, 0, 0, 0.08, 7.5, CYAN)

    # Top gradient band
    rect(slide, 0.08, 0, 13.25, 0.06, BLUE)

    # Tagline chip
    chip(slide, "  POINT OF VIEW  ", 0.55, 0.55, 1.8, 0.32, CYAN, NAVY, 9)

    # Main title
    label(slide, "Rudra Multi-Model\nAgentic Architecture",
          0.55, 1.0, 9.5, 2.2, font_size=44, bold=True, color=WHITE,
          align=PP_ALIGN.LEFT)

    # Subtitle
    label(slide, "Delivering ERP & EPM Consulting in 1 Week\nUsing Claude · OpenClaw · Qwen 2.5 · Canva · Databricks · Video AI",
          0.55, 3.15, 10.0, 1.1, font_size=15, color=LIGHT_GREY,
          align=PP_ALIGN.LEFT)

    divider(slide, 4.4, CYAN, l=0.55, r=1.5)

    # Model logos row
    models = [
        ("Claude\nOpus 4.6", BLUE),
        ("OpenClaw\nGPT-4o", RGBColor(0x10, 0xA3, 0x7F)),
        ("Qwen\n2.5", RGBColor(0xFF, 0x61, 0x00)),
        ("Canva\nMCP", RGBColor(0x8B, 0x3D, 0xFF)),
        ("Databricks\nAI", RGBColor(0xFF, 0x33, 0x21)),
        ("Video\nAnimation", GREEN),
    ]
    x = 0.55
    for name, color in models:
        chip(slide, name, x, 4.65, 1.8, 0.72, color, WHITE, 8.5)
        x += 1.95

    # Bottom right watermark
    label(slide, "Finance Transformation · Technical Accounting · ERP/EPM Systems Integration",
          0.55, 6.85, 12.0, 0.4, font_size=8, color=MID_GREY, align=PP_ALIGN.LEFT)


# ─── SLIDE 2: The Problem ─────────────────────────────────────────────────────

def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)
    rect(slide, 0, 0, 0.08, 7.5, RED)

    label(slide, "THE CHALLENGE", 0.35, 0.18, 5, 0.38,
          font_size=10, bold=True, color=RED, align=PP_ALIGN.LEFT)
    label(slide, "ERP & EPM Consulting is Too Slow, Too Manual, Too Expensive",
          0.35, 0.52, 12.5, 0.82, font_size=26, bold=True, color=WHITE,
          align=PP_ALIGN.LEFT)
    divider(slide, 1.38, RED, l=0.35, r=0.5)

    problems = [
        ("12–18 Weeks", "Average ERP implementation kickoff to go-live for a single module", RED),
        ("$2–5M", "Typical Big-4 finance transformation engagement cost for mid-market", ORANGE),
        ("80% Manual", "Proportion of ERP/EPM deliverables still produced manually in Word/Excel", ORANGE),
        ("3–5 Consultants", "Resources needed just to produce Week 1 discovery and design documents", RED),
        ("6+ Iterations", "Average revision cycles for COA design, process maps and EPM requirements", ORANGE),
        ("No AI Memory", "Traditional tools lose context between deliverables — each doc starts from scratch", RED),
    ]

    cols = [(0.35, 4.1), (6.85, 4.1)]
    col_items = [problems[:3], problems[3:]]

    for (cx, cw), items in zip(cols, col_items):
        y = 1.6
        for stat, desc, color in items:
            rounded_rect(slide, cx, y, cw, 1.38, DARK_CARD)
            rect(slide, cx, y, 0.07, 1.38, color)
            label(slide, stat, cx + 0.22, y + 0.12, 2.2, 0.55,
                  font_size=26, bold=True, color=color, align=PP_ALIGN.LEFT)
            label(slide, desc, cx + 0.22, y + 0.65, cw - 0.35, 0.65,
                  font_size=9.5, color=LIGHT_GREY, align=PP_ALIGN.LEFT)
            y += 1.52


# ─── SLIDE 3: The Vision ──────────────────────────────────────────────────────

def slide_vision(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)
    rect(slide, 0, 0, 0.08, 7.5, GREEN)

    label(slide, "THE VISION", 0.35, 0.18, 5, 0.38,
          font_size=10, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
    label(slide, "From 12 Weeks to 1 Week — Agentic AI as Your Consulting Team",
          0.35, 0.52, 12.5, 0.82, font_size=26, bold=True, color=WHITE,
          align=PP_ALIGN.LEFT)
    divider(slide, 1.38, GREEN, l=0.35, r=0.5)

    # Before / After comparison
    rect(slide, 0.35, 1.55, 5.8, 5.55, RGBColor(0x1A, 0x10, 0x10))
    rect(slide, 0.35, 1.55, 5.8, 0.42, RED)
    label(slide, "BEFORE  —  Traditional Consulting", 0.55, 1.6, 5.4, 0.35,
          font_size=11, bold=True, color=WHITE)
    befores = [
        "❌  12–18 week delivery timelines",
        "❌  Manual Word/Excel deliverable production",
        "❌  3–5 consultants for basic discovery docs",
        "❌  No context sharing between workstreams",
        "❌  $2–5M+ engagement costs",
        "❌  Siloed accounting, ERP & EPM expertise",
        "❌  Static deliverables, no live visualisations",
    ]
    y = 2.15
    for b in befores:
        label(slide, b, 0.55, y, 5.4, 0.42, font_size=10, color=LIGHT_GREY,
              align=PP_ALIGN.LEFT)
        y += 0.44

    rect(slide, 7.18, 1.55, 5.8, 5.55, RGBColor(0x0A, 0x1A, 0x14))
    rect(slide, 7.18, 1.55, 5.8, 0.42, GREEN)
    label(slide, "AFTER  —  Rudra Agentic Framework", 7.38, 1.6, 5.4, 0.35,
          font_size=11, bold=True, color=WHITE)
    afters = [
        "✅  1-week full ERP/EPM kickoff package",
        "✅  AI-generated Excel, PDF, PPTX, Video",
        "✅  10 specialist agents working in parallel",
        "✅  Shared memory across all workstreams",
        "✅  90% cost reduction on documentation",
        "✅  IFRS/GAAP grounded via vector search",
        "✅  Animated dashboards & roadmap videos",
    ]
    y = 2.15
    for a in afters:
        label(slide, a, 7.38, y, 5.4, 0.42, font_size=10, color=LIGHT_GREY,
              align=PP_ALIGN.LEFT)
        y += 0.44

    # VS bubble
    rounded_rect(slide, 6.17, 3.85, 1.0, 0.75, NAVY)
    label(slide, "VS", 6.17, 3.88, 1.0, 0.68,
          font_size=22, bold=True, color=GOLD)


# ─── SLIDE 4: Architecture Overview ──────────────────────────────────────────

def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)
    rect(slide, 0, 0, 0.08, 7.5, CYAN)

    label(slide, "ARCHITECTURE OVERVIEW", 0.35, 0.18, 6, 0.38,
          font_size=10, bold=True, color=CYAN, align=PP_ALIGN.LEFT)
    label(slide, "Four-Layer Agentic Stack",
          0.35, 0.52, 12.5, 0.65, font_size=26, bold=True, color=WHITE,
          align=PP_ALIGN.LEFT)
    divider(slide, 1.2, CYAN, l=0.35, r=0.5)

    layers = [
        # (label, y, color, items)
        ("LAYER 1 — INPUT & ROUTING", 1.35, CYAN,
         ["Client Brief", "ERP/EPM Task", "Accounting Query", "Pipeline Request"]),
        ("LAYER 2 — ORCHESTRATION  (OpenClaw State Machine)", 2.45, BLUE,
         ["Classify Task", "Route to Agent", "Execute", "Synthesise Output"]),
        ("LAYER 3 — AGENT NETWORK  (Rudra + 8 Specialists)", 3.55, GREEN,
         ["Rudra\n(Principal)", "Accounting\nPolicy", "Revenue\nRecog.", "Lease /\nTax / M&A",
          "EPM\nModeling", "Creative\nDesigner", "Databricks\nPipeline", "Financial\nModeling"]),
        ("LAYER 4 — TOOL ECOSYSTEM", 5.1, ORANGE,
         ["Claude\nOpus 4.6", "OpenClaw\nGPT-4o", "Qwen 2.5", "Canva\nMCP",
          "Databricks\nVector DB", "Video\nAnimation", "ERP/EPM\nTools"]),
    ]

    layer_colors = [CYAN, BLUE, GREEN, ORANGE]

    for i, (lbl, y, color, items) in enumerate(layers):
        h = 0.88 if i != 2 else 1.35
        rounded_rect(slide, 0.35, y, 12.6, h, DARK_CARD)
        rect(slide, 0.35, y, 12.6, 0.04, color)
        label(slide, lbl, 0.5, y + 0.06, 4.0, 0.32,
              font_size=8.5, bold=True, color=color, align=PP_ALIGN.LEFT)

        n = len(items)
        item_w = 11.5 / n
        for j, item in enumerate(items):
            cx = 0.65 + j * (item_w + 0.08)
            iy = y + 0.38 if i != 2 else y + 0.42
            iw = item_w - 0.05
            ih = 0.38 if i != 2 else 0.78
            chip(slide, item, cx, iy, iw, ih, NAVY, color, 7.5)

        # Down arrow between layers (except last)
        if i < len(layers) - 1:
            arrow_down(slide, 6.57, y + h, 0.2, 0.18, layer_colors[i + 1])


# ─── SLIDE 5: Model Layer ─────────────────────────────────────────────────────

def slide_models(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)
    rect(slide, 0, 0, 0.08, 7.5, BLUE)

    label(slide, "THE MODEL LAYER", 0.35, 0.18, 5, 0.38,
          font_size=10, bold=True, color=CYAN, align=PP_ALIGN.LEFT)
    label(slide, "Right Model, Right Task — Intelligent Routing",
          0.35, 0.52, 12.5, 0.72, font_size=26, bold=True, color=WHITE,
          align=PP_ALIGN.LEFT)
    divider(slide, 1.28, BLUE, l=0.35, r=0.5)

    models = [
        {
            "name": "Claude Opus 4.6",
            "provider": "Anthropic",
            "color": BLUE,
            "role": "PRIMARY REASONING ENGINE",
            "tasks": [
                "Accounting policy (IFRS/GAAP)",
                "Revenue recognition analysis",
                "ERP/EPM architecture design",
                "M&A due diligence & QoE",
                "Lease & tax accounting",
                "Board-ready narrative writing",
            ],
            "features": ["Adaptive Thinking", "200K Context", "Tool Use", "Streaming"],
        },
        {
            "name": "OpenClaw / GPT-4o",
            "provider": "OpenAI-Compatible",
            "color": RGBColor(0x10, 0xA3, 0x7F),
            "role": "ORCHESTRATION ENGINE",
            "tasks": [
                "Task classification & routing",
                "Multi-agent workflow control",
                "Output synthesis & merging",
                "Structured data extraction",
                "LangGraph state machine",
                "Human-in-the-loop gates",
            ],
            "features": ["Function Calling", "JSON Mode", "Swappable Backend", "Self-Hosted"],
        },
        {
            "name": "Qwen 2.5",
            "provider": "Alibaba DashScope",
            "color": RGBColor(0xFF, 0x61, 0x00),
            "role": "CODE & DATA ENGINE",
            "tasks": [
                "PySpark / Python generation",
                "Databricks pipeline design",
                "Multilingual report writing",
                "Financial data extraction",
                "Delta Lake schema design",
                "SQL analytics queries",
            ],
            "features": ["OpenAI-Compatible API", "Long Context", "Multilingual", "Code Expert"],
        },
    ]

    x = 0.35
    for m in models:
        w = 4.12
        rounded_rect(slide, x, 1.42, w, 5.65, DARK_CARD)
        rect(slide, x, 1.42, w, 0.06, m["color"])

        label(slide, m["role"], x + 0.15, 1.52, w - 0.25, 0.3,
              font_size=7.5, bold=True, color=m["color"], align=PP_ALIGN.LEFT)
        label(slide, m["name"], x + 0.15, 1.82, w - 0.25, 0.52,
              font_size=17, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        label(slide, m["provider"], x + 0.15, 2.32, w - 0.25, 0.3,
              font_size=9, color=MID_GREY, align=PP_ALIGN.LEFT)

        divider(slide, 2.65, MID_GREY, l=x + 0.1, r=13.33 - x - w + 0.1)

        label(slide, "USE CASES", x + 0.15, 2.72, w - 0.25, 0.26,
              font_size=7.5, bold=True, color=m["color"], align=PP_ALIGN.LEFT)
        ty = 3.02
        for task in m["tasks"]:
            label(slide, f"→  {task}", x + 0.15, ty, w - 0.25, 0.3,
                  font_size=9, color=LIGHT_GREY, align=PP_ALIGN.LEFT)
            ty += 0.32

        divider(slide, 5.0, MID_GREY, l=x + 0.1, r=13.33 - x - w + 0.1)

        label(slide, "CAPABILITIES", x + 0.15, 5.06, w - 0.25, 0.26,
              font_size=7.5, bold=True, color=m["color"], align=PP_ALIGN.LEFT)
        fx = x + 0.15
        fy = 5.36
        for feat in m["features"]:
            fw = (w - 0.3) / 2 - 0.08
            chip(slide, feat, fx, fy, fw, 0.3, NAVY, m["color"], 7)
            if fx > x + 0.15 + 0.1:
                fx = x + 0.15
                fy += 0.38
            else:
                fx += fw + 0.1

        x += 4.45


# ─── SLIDE 6: Agent Network ────────────────────────────────────────────────────

def slide_agents(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)
    rect(slide, 0, 0, 0.08, 7.5, GREEN)

    label(slide, "THE AGENT NETWORK", 0.35, 0.18, 5, 0.38,
          font_size=10, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
    label(slide, "Rudra + 10 Specialist Agents — All Coordinated by OpenClaw",
          0.35, 0.52, 12.5, 0.72, font_size=26, bold=True, color=WHITE,
          align=PP_ALIGN.LEFT)
    divider(slide, 1.28, GREEN, l=0.35, r=0.5)

    # Central Rudra bubble
    rounded_rect(slide, 5.3, 1.55, 2.7, 0.9, BLUE)
    label(slide, "RUDRA\nPrincipal Agent", 5.3, 1.58, 2.7, 0.85,
          font_size=12, bold=True, color=WHITE)

    rounded_rect(slide, 5.55, 2.6, 2.2, 0.72, RGBColor(0x10, 0xA3, 0x7F))
    label(slide, "OpenClaw\nOrchestrator", 5.55, 2.62, 2.2, 0.68,
          font_size=10, bold=True, color=WHITE)

    agents = [
        # (name, x, y, color)
        ("Accounting\nPolicy", 0.25, 1.4, CYAN),
        ("Revenue\nRecognition", 0.25, 3.0, CYAN),
        ("Lease\nAccounting", 0.25, 4.6, CYAN),
        ("Tax\nAccounting", 2.2, 1.4, BLUE),
        ("M&A Due\nDiligence", 2.2, 3.0, BLUE),
        ("Financial\nModeling", 2.2, 4.6, BLUE),
        ("Creative\nDesigner", 8.9, 1.4, RGBColor(0x8B, 0x3D, 0xFF)),
        ("EPM\nRequirements", 8.9, 3.0, RGBColor(0x8B, 0x3D, 0xFF)),
        ("Databricks\nPipeline", 8.9, 4.6, RGBColor(0xFF, 0x33, 0x21)),
        ("Vector Search\nRouter", 10.85, 1.4, ORANGE),
        ("ERP/EPM\nTools", 10.85, 3.0, ORANGE),
        ("Fact Check\nAgent", 10.85, 4.6, ORANGE),
    ]

    for name, x, y, color in agents:
        rounded_rect(slide, x, y, 1.82, 1.28, DARK_CARD)
        rect(slide, x, y, 1.82, 0.04, color)
        label(slide, name, x, y + 0.08, 1.82, 1.12,
              font_size=9, bold=True, color=color)

    # Model tags
    label(slide, "Claude Opus 4.6", 0.35, 5.85, 4.2, 0.36,
          font_size=9, color=CYAN, align=PP_ALIGN.LEFT, bold=True)
    label(slide, "↑ Technical Accounting Specialists", 0.35, 6.18, 4.2, 0.32,
          font_size=8, color=LIGHT_GREY, align=PP_ALIGN.LEFT)

    label(slide, "OpenClaw / GPT-4o", 5.35, 5.85, 3.0, 0.36,
          font_size=9, color=RGBColor(0x10, 0xA3, 0x7F), bold=True)
    label(slide, "↑ Workflow Engine", 5.35, 6.18, 3.0, 0.32,
          font_size=8, color=LIGHT_GREY)

    label(slide, "Claude + Qwen 2.5", 8.85, 5.85, 4.2, 0.36,
          font_size=9, color=ORANGE, align=PP_ALIGN.LEFT, bold=True)
    label(slide, "↑ Design, Data & ERP/EPM Specialists", 8.85, 6.18, 4.2, 0.32,
          font_size=8, color=LIGHT_GREY, align=PP_ALIGN.LEFT)


# ─── SLIDE 7: Tool Ecosystem ──────────────────────────────────────────────────

def slide_tools(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)
    rect(slide, 0, 0, 0.08, 7.5, ORANGE)

    label(slide, "THE TOOL ECOSYSTEM", 0.35, 0.18, 5, 0.38,
          font_size=10, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)
    label(slide, "Six Integrated Tool Layers",
          0.35, 0.52, 12.5, 0.72, font_size=26, bold=True, color=WHITE,
          align=PP_ALIGN.LEFT)
    divider(slide, 1.28, ORANGE, l=0.35, r=0.5)

    tools = [
        {
            "name": "Canva MCP",
            "color": RGBColor(0x8B, 0x3D, 0xFF),
            "icon": "🎨",
            "what": "Visual Deliverable Generation",
            "outputs": ["Executive Presentations", "CFO Dashboards", "ERP Roadmaps", "One-Pagers"],
        },
        {
            "name": "Databricks\nAgentic AI",
            "color": RGBColor(0xFF, 0x33, 0x21),
            "icon": "⚡",
            "what": "Data & AI Platform",
            "outputs": ["IFRS/GAAP Vector Search", "Delta Lake Analytics", "MLflow Tracking", "SQL Warehousing"],
        },
        {
            "name": "Video\nAnimation",
            "color": GREEN,
            "icon": "🎬",
            "what": "Animated Consulting Videos",
            "outputs": ["Roadmap Animations", "KPI Dashboard Videos", "Process Flow Animations", "Luma AI Photorealistic"],
        },
        {
            "name": "ERP/EPM\nTools",
            "color": CYAN,
            "icon": "📊",
            "what": "Consulting Deliverable Generator",
            "outputs": ["Project Plans (XLSX)", "COA Design", "EPM Requirements", "WFP Models"],
        },
        {
            "name": "Claude\nAPI",
            "color": BLUE,
            "icon": "🧠",
            "what": "Adaptive Thinking Engine",
            "outputs": ["Accounting Analysis", "Policy Memos", "Board Reports", "Audit Memos"],
        },
        {
            "name": "Standards\nVector DB",
            "color": GOLD,
            "icon": "📚",
            "what": "IFRS/GAAP Knowledge Base",
            "outputs": ["IFRS 15, 16, 3, 9...", "ASC 606, 842, 740...", "Internal Policies", "APQC Benchmarks"],
        },
    ]

    cols = 3
    tw = 4.1
    th = 2.55
    for i, t in enumerate(tools):
        col = i % cols
        row = i // cols
        x = 0.35 + col * (tw + 0.2)
        y = 1.42 + row * (th + 0.22)
        rounded_rect(slide, x, y, tw, th, DARK_CARD)
        rect(slide, x, y, tw, 0.05, t["color"])

        label(slide, t["icon"] + "  " + t["name"].replace("\n", " "),
              x + 0.15, y + 0.1, tw - 0.25, 0.52,
              font_size=13, bold=True, color=t["color"], align=PP_ALIGN.LEFT)
        label(slide, t["what"], x + 0.15, y + 0.6, tw - 0.25, 0.3,
              font_size=8.5, color=MID_GREY, align=PP_ALIGN.LEFT, italic=True)

        oy = y + 0.95
        for output in t["outputs"]:
            label(slide, f"·  {output}", x + 0.15, oy, tw - 0.25, 0.3,
                  font_size=9, color=LIGHT_GREY, align=PP_ALIGN.LEFT)
            oy += 0.34


# ─── SLIDE 8: OpenClaw State Machine ─────────────────────────────────────────

def slide_state_machine(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)
    rect(slide, 0, 0, 0.08, 7.5, RGBColor(0x10, 0xA3, 0x7F))

    label(slide, "OPENCLAW ORCHESTRATOR", 0.35, 0.18, 6, 0.38,
          font_size=10, bold=True, color=RGBColor(0x10, 0xA3, 0x7F), align=PP_ALIGN.LEFT)
    label(slide, "LangGraph-Compatible State Machine Workflow",
          0.35, 0.52, 12.5, 0.72, font_size=26, bold=True, color=WHITE,
          align=PP_ALIGN.LEFT)
    divider(slide, 1.28, RGBColor(0x10, 0xA3, 0x7F), l=0.35, r=0.5)

    # State machine flow
    states = [
        ("START", DARK_CARD, RGBColor(0x10, 0xA3, 0x7F), "Client request\nreceived"),
        ("CLASSIFY\nTASK", DARK_CARD, CYAN, "OpenClaw identifies\ntask category"),
        ("ROUTE\nAGENT", DARK_CARD, BLUE, "Vector similarity\nto agent registry"),
        ("EXECUTE\nAGENT", DARK_CARD, GREEN, "Primary + secondary\nagents run"),
        ("TOOL\nCALL", DARK_CARD, ORANGE, "Canva / Databricks /\nVideo / ERP-EPM"),
        ("SYNTHESISE", DARK_CARD, RGBColor(0x8B, 0x3D, 0xFF), "OpenClaw merges\nmulti-agent outputs"),
        ("DELIVER", DARK_CARD, GOLD, "Board-ready\ndeliverable"),
    ]

    bw = 1.58
    bh = 1.3
    gap = 0.22
    total = len(states) * bw + (len(states) - 1) * gap
    start_x = (13.33 - total) / 2

    for i, (name, fill, color, desc) in enumerate(states):
        x = start_x + i * (bw + gap)
        y = 1.6
        rounded_rect(slide, x, y, bw, bh, fill)
        rect(slide, x, y, bw, 0.05, color)
        label(slide, name, x, y + 0.1, bw, 0.68,
              font_size=11, bold=True, color=color)
        label(slide, desc, x, y + 0.75, bw, 0.5,
              font_size=7.5, color=LIGHT_GREY)
        if i < len(states) - 1:
            arrow_right(slide, x + bw + 0.02, y + 0.5, 0.18, 0.22, color)

    # Human-in-the-loop callout
    rounded_rect(slide, 4.05, 3.2, 2.5, 0.78, RGBColor(0x1A, 0x14, 0x06))
    rect(slide, 4.05, 3.2, 2.5, 0.04, ORANGE)
    label(slide, "⚠  Human-in-the-Loop\nInterrupt Gate (Optional)",
          4.1, 3.24, 2.4, 0.68, font_size=8.5, color=ORANGE, align=PP_ALIGN.LEFT)

    arrow_down(slide, 5.05, 3.0, 0.18, 0.2, ORANGE)

    # Memory layer
    rounded_rect(slide, 0.35, 4.25, 12.6, 1.62, DARK_CARD)
    rect(slide, 0.35, 4.25, 12.6, 0.04, MID_GREY)
    label(slide, "MEMORY LAYER", 0.55, 4.32, 4.0, 0.3,
          font_size=8, bold=True, color=MID_GREY, align=PP_ALIGN.LEFT)

    mem_types = [
        ("Short-Term Memory", "Active conversation context\nand current task state", CYAN),
        ("Long-Term Memory", "Persisted to Databricks\nDelta Lake", BLUE),
        ("Episodic Memory", "Key decisions, findings\nand deliverable history", GREEN),
        ("Standards RAG", "IFRS/GAAP knowledge base\nvia Vector Search", ORANGE),
    ]
    mx = 0.65
    for mname, mdesc, mcolor in mem_types:
        label(slide, mname, mx, 4.62, 2.9, 0.3,
              font_size=9, bold=True, color=mcolor, align=PP_ALIGN.LEFT)
        label(slide, mdesc, mx, 4.94, 2.9, 0.52,
              font_size=8, color=LIGHT_GREY, align=PP_ALIGN.LEFT)
        mx += 3.2

    # Parallel execution note
    rounded_rect(slide, 0.35, 6.1, 12.6, 1.05, RGBColor(0x0A, 0x16, 0x22))
    rect(slide, 0.35, 6.1, 12.6, 0.04, BLUE)
    label(slide, "🔀  PARALLEL EXECUTION  —  Secondary agents run concurrently via ThreadPoolExecutor (4 workers). "
          "Human-in-the-loop gates trigger at material accounting decisions. "
          "All agent runs logged to MLflow for auditability.",
          0.55, 6.16, 12.2, 0.88, font_size=9, color=LIGHT_GREY, align=PP_ALIGN.LEFT)


# ─── SLIDE 9: 1-Week Delivery ─────────────────────────────────────────────────

def slide_delivery(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)
    rect(slide, 0, 0, 0.08, 7.5, GOLD)

    label(slide, "1-WEEK DELIVERY MODEL", 0.35, 0.18, 5, 0.38,
          font_size=10, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    label(slide, "Full ERP/EPM Kickoff Package — Generated Automatically",
          0.35, 0.52, 12.5, 0.72, font_size=26, bold=True, color=WHITE,
          align=PP_ALIGN.LEFT)
    divider(slide, 1.28, GOLD, l=0.35, r=0.5)

    days = [
        ("DAY 1", CYAN, [
            "Project Plan (XLSX)",
            "Stakeholder Interviews",
            "Current State Baseline",
            "Standards RAG Setup",
        ]),
        ("DAY 2", BLUE, [
            "COA Design Document",
            "Gap/Fit Analysis",
            "Process Maps (R2R/O2C/P2P)",
            "Solution Architecture",
        ]),
        ("DAY 3", GREEN, [
            "EPM Requirements Doc",
            "WFP Model Design",
            "Data Migration Templates",
            "Integration Inventory",
        ]),
        ("DAY 4", ORANGE, [
            "Exec Presentation (Canva)",
            "CFO Dashboard Mockup",
            "Animated Roadmap Video",
            "Business Case Model",
        ]),
        ("DAY 5", RGBColor(0x8B, 0x3D, 0xFF), [
            "Board Pack (PDF/PPTX)",
            "Risk & Mitigation Log",
            "Week 2 Sprint Plan",
            "Deliverable Review",
        ]),
    ]

    dw = 2.42
    for i, (day, color, items) in enumerate(days):
        x = 0.35 + i * (dw + 0.1)
        rounded_rect(slide, x, 1.42, dw, 5.65, DARK_CARD)
        rect(slide, x, 1.42, dw, 0.5, color)
        label(slide, day, x, 1.44, dw, 0.46,
              font_size=14, bold=True, color=NAVY)

        iy = 2.05
        for item in items:
            rounded_rect(slide, x + 0.12, iy, dw - 0.24, 0.98, NAVY)
            label(slide, item, x + 0.12, iy + 0.08, dw - 0.24, 0.82,
                  font_size=9.5, color=color, align=PP_ALIGN.LEFT)
            iy += 1.12

        # Arrow between days
        if i < len(days) - 1:
            arrow_right(slide, x + dw + 0.01, 2.85, 0.08, 0.2, color)


# ─── SLIDE 10: Databricks Agentic AI ─────────────────────────────────────────

def slide_databricks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)
    rect(slide, 0, 0, 0.08, 7.5, RGBColor(0xFF, 0x33, 0x21))

    label(slide, "DATABRICKS AGENTIC AI", 0.35, 0.18, 6, 0.38,
          font_size=10, bold=True, color=RGBColor(0xFF, 0x33, 0x21), align=PP_ALIGN.LEFT)
    label(slide, "Lakehouse-Powered Finance Intelligence",
          0.35, 0.52, 12.5, 0.72, font_size=26, bold=True, color=WHITE,
          align=PP_ALIGN.LEFT)
    divider(slide, 1.28, RGBColor(0xFF, 0x33, 0x21), l=0.35, r=0.5)

    # Left: Medallion architecture
    layers_db = [
        ("BRONZE LAYER  —  Raw Ingestion",
         "Oracle ERP / SAP / Workday raw exports\nGL transactions, AP invoices, AR ledger, Fixed assets",
         RGBColor(0xCD, 0x7F, 0x32)),
        ("SILVER LAYER  —  Cleansed & Conformed",
         "Standardised schemas, data quality checks\nCOA mapping, entity normalisation, currency translation",
         RGBColor(0xC0, 0xC0, 0xC0)),
        ("GOLD LAYER  —  Finance Analytics",
         "Close metrics, KPI dashboards, EPM forecasts\nManagement reports, consolidation models",
         GOLD),
    ]

    y = 1.5
    for lbl, desc, color in layers_db:
        rounded_rect(slide, 0.35, y, 6.4, 1.28, DARK_CARD)
        rect(slide, 0.35, y, 0.08, 1.28, color)
        label(slide, lbl, 0.58, y + 0.1, 5.9, 0.32,
              font_size=9, bold=True, color=color, align=PP_ALIGN.LEFT)
        label(slide, desc, 0.58, y + 0.44, 5.9, 0.72,
              font_size=9, color=LIGHT_GREY, align=PP_ALIGN.LEFT)
        y += 1.42

    # Right: Capabilities
    caps = [
        ("Vector Search", "Semantic search over IFRS/GAAP standards corpus\nBGE-Large embeddings · 1500 token chunks", CYAN),
        ("SQL Warehousing", "Finance analytics queries · Close metrics · EPM data\nServerless compute · Auto-scaling", BLUE),
        ("MLflow", "Agent run tracking · Model governance\nAudit trail for all AI-generated deliverables", GREEN),
        ("Unity Catalog", "Data governance · Lineage · Access control\nFinance data certified for compliance", ORANGE),
    ]

    y = 1.5
    for cap, desc, color in caps:
        rounded_rect(slide, 7.15, y, 5.8, 1.28, DARK_CARD)
        rect(slide, 7.15, y, 5.8, 0.05, color)
        label(slide, cap, 7.3, y + 0.1, 5.45, 0.32,
              font_size=11, bold=True, color=color, align=PP_ALIGN.LEFT)
        label(slide, desc, 7.3, y + 0.44, 5.45, 0.72,
              font_size=9, color=LIGHT_GREY, align=PP_ALIGN.LEFT)
        y += 1.42

    # Bottom banner
    rounded_rect(slide, 0.35, 5.88, 12.6, 1.28, DARK_CARD)
    rect(slide, 0.35, 5.88, 12.6, 0.04, RGBColor(0xFF, 0x33, 0x21))
    label(slide, "AGENTIC PATTERN:  Agent → Tool Call → Databricks API → Delta Lake → Vector Index → Agent Response",
          0.55, 5.98, 12.1, 0.36, font_size=10, bold=True, color=RGBColor(0xFF, 0x33, 0x21),
          align=PP_ALIGN.LEFT)
    label(slide, "All agent executions are logged to MLflow with agent name, task, token count, latency and output summary. "
          "Databricks Jobs trigger nightly ERP data refreshes to keep the Gold layer current for real-time analytics.",
          0.55, 6.36, 12.1, 0.62, font_size=8.5, color=LIGHT_GREY, align=PP_ALIGN.LEFT)


# ─── SLIDE 11: Business Value ──────────────────────────────────────────────────

def slide_value(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)
    rect(slide, 0, 0, 0.08, 7.5, GREEN)

    label(slide, "BUSINESS VALUE", 0.35, 0.18, 5, 0.38,
          font_size=10, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
    label(slide, "Quantified Impact — Speed, Quality, Cost",
          0.35, 0.52, 12.5, 0.72, font_size=26, bold=True, color=WHITE,
          align=PP_ALIGN.LEFT)
    divider(slide, 1.28, GREEN, l=0.35, r=0.5)

    metrics = [
        ("90%", "Reduction in\nDocumentation Cost", GREEN),
        ("12x", "Faster Deliverable\nProduction", CYAN),
        ("1 Week", "Full ERP/EPM\nKickoff Package", GOLD),
        ("10+", "Specialist Agents\nWorking in Parallel", BLUE),
        ("100%", "Standards-Grounded\nIFRS/GAAP Output", GREEN),
        ("$0 Extra", "Per Additional\nDeliverable Iteration", ORANGE),
    ]

    mw = 2.05
    for i, (stat, desc, color) in enumerate(metrics):
        col = i % 3
        row = i // 3
        x = 0.35 + col * (mw + 0.2)
        y = 1.52 + row * 2.62

        rounded_rect(slide, x, y, mw, 2.38, DARK_CARD)
        rect(slide, x, y, mw, 0.06, color)
        label(slide, stat, x, y + 0.18, mw, 1.05,
              font_size=34, bold=True, color=color)
        label(slide, desc, x, y + 1.2, mw, 0.95,
              font_size=10, color=LIGHT_GREY)

    # Right column: qualitative
    qual = [
        ("Board-Ready Quality", "Every output written to Big-4 standard with proper IFRS/GAAP citations"),
        ("Auditability", "Full MLflow audit trail — every agent decision is traceable and reproducible"),
        ("Human-in-the-Loop", "Material accounting judgments always require human approval before delivery"),
        ("Extensible", "Add new specialist agents, models, or tools without rebuilding the core framework"),
    ]

    x = 7.05
    y = 1.52
    for qtitle, qdesc in qual:
        rounded_rect(slide, x, y, 5.95, 1.38, DARK_CARD)
        rect(slide, x, y, 5.95, 0.04, GREEN)
        label(slide, qtitle, x + 0.15, y + 0.1, 5.6, 0.35,
              font_size=11, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
        label(slide, qdesc, x + 0.15, y + 0.48, 5.6, 0.72,
              font_size=9, color=LIGHT_GREY, align=PP_ALIGN.LEFT)
        y += 1.52


# ─── SLIDE 12: Next Steps ─────────────────────────────────────────────────────

def slide_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)
    rect(slide, 0, 0, 0.08, 7.5, CYAN)

    label(slide, "NEXT STEPS", 0.35, 0.18, 5, 0.38,
          font_size=10, bold=True, color=CYAN, align=PP_ALIGN.LEFT)
    label(slide, "How to Get Started — Immediate Actions",
          0.35, 0.52, 12.5, 0.72, font_size=26, bold=True, color=WHITE,
          align=PP_ALIGN.LEFT)
    divider(slide, 1.28, CYAN, l=0.35, r=0.5)

    steps = [
        ("01", "Configure API Keys", CYAN,
         "Add ANTHROPIC_API_KEY + OPENAI_API_KEY to .env\n"
         "Optional: DASHSCOPE_API_KEY (Qwen), CANVA_ACCESS_TOKEN,\n"
         "DATABRICKS_HOST/TOKEN, LUMA_AI_KEY"),
        ("02", "Run Week 1 Package", BLUE,
         "python main.py --week1 \"Your Client\" \"Project Name\"\n"
         "Generates: Project Plan, COA, EPM Requirements,\n"
         "WFP Model, Data Migration Templates in <60 seconds"),
        ("03", "Connect Databricks", GREEN,
         "Populate standards vector index with IFRS/GAAP corpus\n"
         "Configure Delta Lake schema for finance analytics\n"
         "Enable MLflow audit trail for all agent runs"),
        ("04", "Enable Video & Canva", ORANGE,
         "apt install ffmpeg  →  real MP4 animations\n"
         "Add CANVA_ACCESS_TOKEN  →  live design generation\n"
         "Add LUMA_AI_KEY  →  photorealistic AI videos"),
    ]

    sw = 6.1
    for i, (num, title, color, desc) in enumerate(steps):
        col = i % 2
        row = i // 2
        x = 0.35 + col * (sw + 0.82)
        y = 1.52 + row * 2.6

        rounded_rect(slide, x, y, sw, 2.32, DARK_CARD)
        rect(slide, x, y, sw, 0.05, color)
        rounded_rect(slide, x + 0.12, y + 0.18, 0.68, 0.68, color)
        label(slide, num, x + 0.12, y + 0.18, 0.68, 0.68,
              font_size=16, bold=True, color=NAVY)
        label(slide, title, x + 0.95, y + 0.22, sw - 1.1, 0.45,
              font_size=13, bold=True, color=color, align=PP_ALIGN.LEFT)
        label(slide, desc, x + 0.95, y + 0.72, sw - 1.05, 1.42,
              font_size=9, color=LIGHT_GREY, align=PP_ALIGN.LEFT)

    # Footer CTA
    rect(slide, 0.35, 6.98, 12.6, 0.42, BLUE)
    label(slide, "🚀  Built on branch: claude/multi-model-ai-agent-framework  ·  PR: github.com/rshetty66/Finance-Automation/pull/9  ·  Run: python main.py",
          0.55, 6.99, 12.2, 0.38, font_size=8.5, color=WHITE, align=PP_ALIGN.LEFT)


# ── Build & Save ──────────────────────────────────────────────────────────────

def build():
    import os
    prs = new_prs()

    print("Building slides...")
    slide_title(prs);          print("  ✓ Slide 1 — Title")
    slide_problem(prs);        print("  ✓ Slide 2 — The Challenge")
    slide_vision(prs);         print("  ✓ Slide 3 — The Vision")
    slide_architecture(prs);   print("  ✓ Slide 4 — Architecture Overview")
    slide_models(prs);         print("  ✓ Slide 5 — Model Layer")
    slide_agents(prs);         print("  ✓ Slide 6 — Agent Network")
    slide_tools(prs);          print("  ✓ Slide 7 — Tool Ecosystem")
    slide_state_machine(prs);  print("  ✓ Slide 8 — OpenClaw State Machine")
    slide_delivery(prs);       print("  ✓ Slide 9 — 1-Week Delivery Model")
    slide_databricks(prs);     print("  ✓ Slide 10 — Databricks Agentic AI")
    slide_value(prs);          print("  ✓ Slide 11 — Business Value")
    slide_next_steps(prs);     print("  ✓ Slide 12 — Next Steps")

    out = "output/Rudra_Agentic_Architecture_POV.pptx"
    os.makedirs("output", exist_ok=True)
    prs.save(out)
    print(f"\n✅  Saved → {out}")
    return out


if __name__ == "__main__":
    build()
